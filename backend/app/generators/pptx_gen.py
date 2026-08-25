import os
import hashlib
from pathlib import Path
from datetime import datetime, timezone
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from app.config import settings

class PptxGenerator:
    """
    Sovereign Presentation Generator.
    Produces polished 3-slide executive board briefings and turnaround reviews.
    """
    @staticmethod
    def create_executive_deck(
        data: Optional[Dict[str, Any]] = None,
        output_filename: Optional[str] = None
    ) -> Dict[str, Any]:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # -------------------------------------------------------------
        # SLIDE 1: Title Slide (Executive Industrial Theme)
        # -------------------------------------------------------------
        slide1 = prs.slides.add_slide(blank_layout)
        
        # Background Top Accent Bar
        top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(13, 148, 136) # Teal
        top_bar.line.fill.background()

        # Title container
        tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "MANGALORE REFINERY AND PETROCHEMICALS LIMITED"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(13, 148, 136)

        p1 = tf.add_paragraph()
        p1.text = "EXECUTIVE TURNAROUND & ASSET INTEGRITY BRIEFING"
        p1.font.size = Pt(32)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(15, 23, 42)
        p1.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = "Crude Distillation Unit (CDU-1) Heat Exchanger 11-HX-401 Root Cause Analysis & Mitigation"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(71, 85, 105)
        p2.space_before = Pt(10)

        # Bottom Metadata Card
        meta_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2))
        meta_box.fill.solid()
        meta_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
        meta_box.line.color.rgb = RGBColor(203, 213, 225)
        mtf = meta_box.text_frame
        mtf.word_wrap = True
        mp = mtf.paragraphs[0]
        mp.text = f"SovereignAI Air-Gapped Workbench | Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d')} | Originator: Mechanical Integrity Division"
        mp.font.size = Pt(11)
        mp.font.color.rgb = RGBColor(100, 116, 139)
        mp2 = mtf.add_paragraph()
        mp2.text = "Classification: CONFIDENTIAL MRPL INTERNAL MEMO — Zero Cloud Telemetry Transmitted"
        mp2.font.size = Pt(10)
        mp2.font.bold = True
        mp2.font.color.rgb = RGBColor(185, 28, 28)

        # -------------------------------------------------------------
        # SLIDE 2: NDT Findings & Step-by-Step Calculations
        # -------------------------------------------------------------
        slide2 = prs.slides.add_slide(blank_layout)
        
        # Header
        h_box = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        hp = h_box.text_frame.paragraphs[0]
        hp.text = "1. Critical Ultrasonic NDT Findings & Calculation Derivations"
        hp.font.size = Pt(24)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(15, 23, 42)

        # Left Card: Measured Metrics
        card_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
        card_left.fill.solid()
        card_left.fill.fore_color.rgb = RGBColor(254, 242, 242) # Light red tint
        card_left.line.color.rgb = RGBColor(252, 165, 165)
        cl_tf = card_left.text_frame
        cl_tf.word_wrap = True
        cl_p0 = cl_tf.paragraphs[0]
        cl_p0.text = "DEFECT DISCOVERY & SOP BREACH"
        cl_p0.font.size = Pt(14)
        cl_p0.font.bold = True
        cl_p0.font.color.rgb = RGBColor(185, 28, 28)

        left_points = [
            "Equipment Tag: 11-HX-401 (Crude Preheat Train Pass 2)",
            "Nominal Design Wall Thickness: 5.00 mm",
            "Measured Minimum Wall Thickness: 3.18 mm",
            "SOP-08 Mandatory Cut-Off Limit: 3.50 mm",
            "Safety Margin Delta: -0.32 mm (CRITICAL BREACH)",
            "Corrosion Rate: 0.95 mm/year (Chloride Pitting Attack)",
            "Calculated Remaining Useful Life: EXPIRED (< 45 Days)"
        ]
        for pt in left_points:
            p = cl_tf.add_paragraph()
            p.text = f"• {pt}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.space_before = Pt(6)

        # Right Card: Mathematical Formulas
        card_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
        card_right.fill.solid()
        card_right.fill.fore_color.rgb = RGBColor(240, 253, 250) # Teal tint
        card_right.line.color.rgb = RGBColor(153, 246, 228)
        cr_tf = card_right.text_frame
        cr_tf.word_wrap = True
        cr_p0 = cr_tf.paragraphs[0]
        cr_p0.text = "ENGINEERING FORMULA PROOFS"
        cr_p0.font.size = Pt(14)
        cr_p0.font.bold = True
        cr_p0.font.color.rgb = RGBColor(13, 148, 136)

        right_points = [
            "Formula 1: Thickness Loss (Δt) = 5.00mm - 3.18mm = 1.82 mm",
            "Formula 2: Corrosion Rate (CR) = 1.82mm / 1.92yr = 0.948 mm/yr",
            "Formula 3: SOP Margin = 3.18mm - 3.50mm = -0.32 mm",
            "Formula 4: RUL = (3.18 - 3.50) / 0.948 = -0.338 Years",
            "SOP-08 Section 4.2 Mandate: Direct Category-A emergency isolation required when measured wall falls below 3.50 mm.",
            "Turnaround Time Comparison: Manual Drafting (4.5 hrs) vs SovereignAI (5.6s) -> 99.9% Faster"
        ]
        for pt in right_points:
            p = cr_tf.add_paragraph()
            p.text = f"✔ {pt}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.space_before = Pt(6)

        # -------------------------------------------------------------
        # SLIDE 3: Action Plan & Management Sign-Off
        # -------------------------------------------------------------
        slide3 = prs.slides.add_slide(blank_layout)
        h3_box = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        h3p = h3_box.text_frame.paragraphs[0]
        h3p.text = "2. Turnaround Action Plan & Financial Authorization"
        h3p.font.size = Pt(24)
        h3p.font.bold = True
        h3p.font.color.rgb = RGBColor(15, 23, 42)

        # Action plan table
        actions = [
            ("1", "Immediate Safety Isolation", "Engage bypass line 6\"-BPS-108-A1A to isolate Pass 2 bundle.", "Unit In-Charge", "Immediate (< 24h)"),
            ("2", "Emergency Retubing Order", "Issue SAP Work Order WO-2026-08-4101 for tube bundle procurement.", "Procurement Div.", "3 Days"),
            ("3", "Material Upgrade", "Upgrade replacement bundle to Inconel-625 alloy (prevents chloride stress).", "Integrity Dept.", "Turnaround Window"),
            ("4", "Financial Authorization", "Approve OPEX allocation INR 18.5 Lakhs for turnkey tube pull & replacement.", "CGM (Operations)", "Pending Sign-Off")
        ]

        table_shape = slide3.shapes.add_table(rows=5, cols=5, left=Inches(0.8), top=Inches(1.6), width=Inches(11.7), height=Inches(3.8))
        tbl = table_shape.table
        headers = ["Step", "Action Item", "Engineering Description", "Responsibility", "Timeline"]
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        for i, row_data in enumerate(actions, 1):
            for j, val in enumerate(row_data):
                cell = tbl.cell(i, j)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if i % 2 == 1 else RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)

        if not output_filename:
            output_filename = f"MRPL_Executive_Deck_{int(datetime.now(timezone.utc).timestamp())}.pptx"

        file_path = settings.GENERATED_DIR / output_filename
        prs.save(str(file_path))

        file_size = os.path.getsize(str(file_path))
        with open(str(file_path), "rb") as f:
            file_hash = hashlib.sha256(f.read()).hexdigest()

        return {
            "filename": output_filename,
            "file_type": "PPTX",
            "file_size_bytes": file_size,
            "storage_path": str(file_path),
            "integrity_sha256": file_hash
        }

pptx_generator = PptxGenerator()
